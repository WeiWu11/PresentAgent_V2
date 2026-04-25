import asyncio
import functools
import hashlib
import importlib
import json
import os
import re
import sys
import traceback
import uuid
import subprocess
import tempfile
from contextlib import asynccontextmanager
from copy import deepcopy
from datetime import datetime
from typing import Optional

from pdf2image import convert_from_path
from pptx import Presentation as PptxPresentation

from fastapi import (
    FastAPI,
    File,
    Form,
    HTTPException,
    Request,
    UploadFile,
    WebSocket,
    WebSocketDisconnect,
)
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import FileResponse
sys.path.append("../")
import pptagent.induct as induct
import pptagent.pptgen as pptgen
from pptagent.document import Document
from pptagent.model_utils import ModelManager, parse_pdf
from pptagent.multimodal import ImageLabler
from pptagent.presentation import Presentation
from pptagent.research import DeepResearchAdapter, DocumentMediaResearcher, ResearchDossier
from pptagent.utils import Config, get_logger, package_join, pjoin, ppt_to_images_async


async def run_blocking(func, *args, **kw):
    loop = asyncio.get_running_loop()
    return await loop.run_in_executor(None, functools.partial(func, *args, **kw))

async def run_cmd(cmd: list[str]):
    proc = await asyncio.create_subprocess_exec(
        *cmd, stdout=asyncio.subprocess.PIPE, stderr=asyncio.subprocess.PIPE
    )
    stdout, stderr = await proc.communicate()
    if proc.returncode != 0:
        raise RuntimeError(f"{' '.join(cmd)}\n{stderr.decode()}")
    return stdout
# ----------------------------------------------------------------------

# constants
DEBUG = True if len(sys.argv) == 1 else False
RUNS_DIR = package_join("runs")
DOCUMENT_STAGES = [
    "PPT Parsing",
    "PDF Parsing",
    "Document Refinement",
    "External Media Retrieval",
    "PPT Analysis",
    "PPT Generation",
    "Success!",
]
TOPIC_STAGES = [
    "PPT Parsing",
    "Deep Research",
    "Dossier Adaptation",
    "PPT Analysis",
    "PPT Generation",
    "Success!",
]


ppt_video_progress_store: dict[str, dict] = {}

models = ModelManager()
DOC_MEDIA_MAX_CONCEPTS = 3
DOC_MEDIA_MAX_WAIT_SECONDS = 600.0
DOC_MEDIA_POLL_INTERVAL_SECONDS = 60.0
DOC_MEDIA_DEEPRESEARCH_CONDA_ENV = "react_infer_env"
DOC_MEDIA_DEEPRESEARCH_CONDA_EXE = "conda"


@asynccontextmanager
async def lifespan(_: FastAPI):
    assert await models.test_connections(), "Model connection test failed"
    yield

# server
logger = get_logger(__name__)
app = FastAPI(lifespan=lifespan)
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)
progress_store: dict[str, dict] = {}
active_connections: dict[str, WebSocket] = {}


ppt_video_active_connections: dict[str, WebSocket] = {}

SPEAKER_LINE_RE = re.compile(r"^\s*Speaker\s+([AB])\s*:\s*(.+?)\s*$", re.I)


class ProgressManager:
    def __init__(self, task_id: str, stages: list[str], debug: bool = True):
        self.task_id = task_id
        self.stages = stages
        self.debug = debug
        self.failed = False
        self.current_stage = 0
        self.total_stages = len(stages)

    async def report_progress(self):
        assert self.task_id in active_connections, "WebSocket connection closed"
        self.current_stage += 1
        progress = int((self.current_stage / self.total_stages) * 100)
        await send_progress(
            active_connections[self.task_id],
            f"Stage: {self.stages[self.current_stage - 1]}",
            progress,
        )

    async def fail_stage(self, error_message: str):
        await send_progress(
            active_connections[self.task_id],
            f"{self.stages[self.current_stage]} Error: {error_message}",
            100,
        )
        self.failed = True
        active_connections.pop(self.task_id, None)
        if self.debug:
            logger.error(
                f"{self.task_id}: {self.stages[self.current_stage]} Error: {error_message}"
            )

@app.post("/api/upload")
async def create_task(
        pptxFile: UploadFile = File(None),
        pdfFile: UploadFile = File(None),
        topic: str = Form(None),
        numberOfPages: int = Form(...),
):
    task_id = datetime.now().strftime("20%y-%m-%d") + "/" + str(uuid.uuid4())
    logger.info(f"task created: {task_id}")
    os.makedirs(pjoin(RUNS_DIR, task_id))
    task = {
        "numberOfPages": numberOfPages,
        "pptx": "default_template",
    }
    if pptxFile is not None:
        pptx_blob = await pptxFile.read()
        pptx_md5 = hashlib.md5(pptx_blob).hexdigest()
        task["pptx"] = pptx_md5
        pptx_dir = pjoin(RUNS_DIR, "pptx", pptx_md5)
        if not os.path.exists(pptx_dir):
            os.makedirs(pptx_dir, exist_ok=True)
            with open(pjoin(pptx_dir, "source.pptx"), "wb") as f:
                f.write(pptx_blob)
    if pdfFile is not None:
        pdf_blob = await pdfFile.read()
        pdf_md5 = hashlib.md5(pdf_blob).hexdigest()
        task["pdf"] = pdf_md5
        pdf_dir = pjoin(RUNS_DIR, "pdf", pdf_md5)
        if not os.path.exists(pdf_dir):
            os.makedirs(pdf_dir, exist_ok=True)
            with open(pjoin(pdf_dir, "source.pdf"), "wb") as f:
                f.write(pdf_blob)
    if topic is not None:
        task["topic"] = topic.strip()
    progress_store[task_id] = task
    # Start the PPT generation task asynchronously
    asyncio.create_task(ppt_gen(task_id))
    return {"task_id": task_id.replace("/", "|")}


async def send_progress(websocket: Optional[WebSocket], status: str, progress: int):
    if websocket is None:
        logger.info(f"websocket is None, status: {status}, progress: {progress}")
        return
    await websocket.send_json({"progress": progress, "status": status})


async def send_ppt_video_progress(task_id: str):
    if task_id in ppt_video_active_connections:
        progress_data = ppt_video_progress_store[task_id]
        await ppt_video_active_connections[task_id].send_json(progress_data)
    else:
        logger.warning(f"No WebSocket connection for PPT video task {task_id}")


@app.websocket("/wsapi/{task_id}")
async def websocket_endpoint(websocket: WebSocket, task_id: str):
    task_id = task_id.replace("|", "/")
    if task_id in progress_store:
        await websocket.accept()
    else:
        raise HTTPException(status_code=404, detail="Task not found")
    active_connections[task_id] = websocket
    try:
        while True:
            await websocket.receive_text()
    except WebSocketDisconnect:
        logger.info("websocket disconnected: %s", task_id)
        active_connections.pop(task_id, None)


@app.get("/api/download")
async def download(task_id: str):
    task_id = task_id.replace("|", "/")
    if not os.path.exists(pjoin(RUNS_DIR, task_id)):
        raise HTTPException(status_code=404, detail="Task not created yet")
    file_path = pjoin(RUNS_DIR, task_id, "final.pptx")
    if os.path.exists(file_path):
        return FileResponse(
            file_path,
            media_type="application/pptx",
            headers={"Content-Disposition": "attachment; filename=pptagent.pptx"},
        )
    raise HTTPException(status_code=404, detail="Task not finished yet")


@app.post("/api/feedback")
async def feedback(request: Request):
    body = await request.json()
    feedback = body.get("feedback")
    task_id = body.get("task_id")

    with open(pjoin(RUNS_DIR, "feedback", f"{task_id}.txt"), "w") as f:
        f.write(feedback)
    return {"message": "Feedback submitted successfully"}


@app.get("/")
async def hello():
    return {"message": "Hello, World!"}


@app.post("/api/ppt-to-video")
async def create_ppt_video_task(pptFile: UploadFile = File(...)):
    task_id = str(uuid.uuid4())
    logger.info(f"PPT2Presentation task created: {task_id}")

    task_dir = pjoin(RUNS_DIR, "ppt_video", task_id)
    os.makedirs(task_dir, exist_ok=True)

    ppt_blob = await pptFile.read()
    ppt_path = pjoin(task_dir, "source.pptx")
    with open(ppt_path, "wb") as f:
        f.write(ppt_blob)

    ppt_video_progress_store[task_id] = {
        "status": "processing",
        "current_step": 1,
        "current_slide": 0,
        "total_slides": 0,
        "progress_percentage": 0,
        "task_dir": task_dir,
        "ppt_path": ppt_path
    }

    asyncio.create_task(process_ppt_to_video(task_id))
    return {"task_id": task_id}


@app.websocket("/wsapi/ppt-to-video/{task_id}")
async def websocket_ppt_video_endpoint(websocket: WebSocket, task_id: str):
    if task_id in ppt_video_progress_store:
        await websocket.accept()
    else:
        raise HTTPException(status_code=404, detail="Task not found")
    ppt_video_active_connections[task_id] = websocket
    try:
        while True:
            await websocket.receive_text()
    except WebSocketDisconnect:
        logger.info("PPT video websocket disconnected: %s", task_id)
        ppt_video_active_connections.pop(task_id, None)


async def process_ppt_to_video(task_id: str):
    task_dir = ppt_video_progress_store[task_id]["task_dir"]
    try:
        ppt_path = ppt_video_progress_store[task_id]["ppt_path"]
        ppt_video_progress_store[task_id].update(current_step=1, progress_percentage=10.00)
        await send_ppt_video_progress(task_id)

        pdf_path = pjoin(task_dir, "source.pdf")
        await run_cmd([
            "libreoffice", "--headless", "--convert-to", "pdf",
            ppt_path, "--outdir", task_dir
        ])

        images_from_path = await run_blocking(convert_from_path, pdf_path)
        prs = await run_blocking(PptxPresentation, ppt_path)

        if len(images_from_path) != len(prs.slides):
            raise Exception("PPT页数与生成的图片数量不匹配")

        ppt_video_progress_store[task_id].update(
            total_slides=len(prs.slides), progress_percentage=20.00
        )
        await send_ppt_video_progress(task_id)

        ppt_video_progress_store[task_id].update(current_step=2, progress_percentage=30.00)
        await send_ppt_video_progress(task_id)

        video_segments = []
        with tempfile.TemporaryDirectory() as temp_path:
            for i, (slide, image) in enumerate(zip(prs.slides, images_from_path)):
                ppt_video_progress_store[task_id].update(
                    current_slide=i + 1,
                    progress_percentage=round(30 + (i / len(prs.slides)) * 40,2)
                )
                await send_ppt_video_progress(task_id)

                notes = slide.notes_slide.notes_text_frame.text if slide.has_notes_slide else ""
                if not notes.strip():
                    notes = f"This is the {i + 1} page"

                image_path = pjoin(temp_path, f"frame_{i}.jpg")
                image.save(image_path)

                audio_path = pjoin(temp_path, f"frame_{i}.wav")
                await generate_tts_audio(notes, audio_path)

                video_segment_path = await create_video_segment(
                    image_path, audio_path, temp_path, i
                )
                video_segments.append(video_segment_path)

            ppt_video_progress_store[task_id].update(current_step=3, progress_percentage=80.00)
            await send_ppt_video_progress(task_id)

            output_video_path = pjoin(task_dir, "output.mp4")
            await merge_video_segments(video_segments, output_video_path)

        ppt_video_progress_store[task_id].update(
            status="completed",
            progress_percentage=100.00,
            video_url=f"/api/ppt-to-video/download/{task_id}",
        )
        await send_ppt_video_progress(task_id)

    except Exception as e:
        logger.error(f"PPT2Presentation task failed {task_id}: {e}")
        ppt_video_progress_store[task_id].update(status="failed", error_message=str(e))
        await send_ppt_video_progress(task_id)


async def generate_tts_audio(text: str, output_path: str):
    try:
        sys.path.append(pjoin(os.path.dirname(__file__), "MegaTTS3"))
        from tts.infer_cli import MegaTTS3DiTInfer
        from tts.utils.audio_utils.io import save_wav

        infer = MegaTTS3DiTInfer(ckpt_root=pjoin(os.path.dirname(__file__), "MegaTTS3", "checkpoints"))
        assets_dir = pjoin(os.path.dirname(__file__), "MegaTTS3", "assets")
        speaker_a_prompt = os.getenv(
            "TTS_SPEAKER_A_PROMPT",
            pjoin(assets_dir, "English_prompt.wav"),
        )
        speaker_b_prompt = os.getenv(
            "TTS_SPEAKER_B_PROMPT",
            pjoin(assets_dir, "Chinese_prompt.wav"),
        )

        def parse_speaker_segments(script: str) -> list[tuple[str, str]]:
            segments: list[tuple[str, str]] = []
            current_speaker = "A"
            buffer: list[str] = []
            for raw_line in script.splitlines():
                line = raw_line.strip()
                if not line:
                    continue
                match = SPEAKER_LINE_RE.match(line)
                if match:
                    if buffer:
                        segments.append((current_speaker, " ".join(buffer).strip()))
                        buffer = []
                    current_speaker = match.group(1).upper()
                    content = match.group(2).strip()
                    if content:
                        buffer.append(content)
                else:
                    buffer.append(line)
            if buffer:
                segments.append((current_speaker, " ".join(buffer).strip()))
            if not segments:
                segments.append(("A", script.strip()))
            return [(speaker, content) for speaker, content in segments if content]

        def build_resource_context(prompt_audio_path: str):
            with open(prompt_audio_path, "rb") as f:
                audio_bytes = f.read()
            latent_file = None
            potential_npy = os.path.splitext(prompt_audio_path)[0] + ".npy"
            if os.path.isfile(potential_npy):
                latent_file = potential_npy
            return infer.preprocess(audio_bytes, latent_file)

        def synthesize_segments(script: str, destination: str):
            segments = parse_speaker_segments(script)
            speaker_contexts = {
                "A": build_resource_context(speaker_a_prompt),
                "B": build_resource_context(speaker_b_prompt),
            }
            with tempfile.TemporaryDirectory() as temp_dir:
                segment_paths: list[str] = []
                for idx, (speaker, content) in enumerate(segments):
                    wav_bytes = infer.forward(
                        speaker_contexts[speaker],
                        content,
                        time_step=32,
                        p_w=1.6,
                        t_w=2.5,
                    )
                    segment_path = pjoin(temp_dir, f"speaker_{speaker}_{idx}.wav")
                    save_wav(wav_bytes, segment_path)
                    segment_paths.append(segment_path)
                if len(segment_paths) == 1:
                    shutil.copyfile(segment_paths[0], destination)
                    return
                list_path = pjoin(temp_dir, "segments.txt")
                with open(list_path, "w", encoding="utf-8") as f:
                    for segment_path in segment_paths:
                        f.write(f"file '{segment_path}'\n")
                subprocess.run(
                    [
                        "ffmpeg",
                        "-y",
                        "-f",
                        "concat",
                        "-safe",
                        "0",
                        "-i",
                        list_path,
                        "-c",
                        "copy",
                        destination,
                    ],
                    check=True,
                    stdout=subprocess.DEVNULL,
                    stderr=subprocess.DEVNULL,
                )

        synthesize_segments(text, output_path)

    except Exception as e:
        logger.error(f"TTS failed: {str(e)}")
        import numpy as np
        import wave

        sample_rate = 22050
        duration = 3.0
        samples = np.zeros(int(sample_rate * duration), dtype=np.int16)

        with wave.open(output_path, 'w') as wav_file:
            wav_file.setnchannels(1)
            wav_file.setsampwidth(2)
            wav_file.setframerate(sample_rate)
            wav_file.writeframes(samples.tobytes())


async def create_video_segment(image_path: str, audio_path: str, temp_path: str, index: int):
    output_path = pjoin(temp_path, f"segment_{index}.mp4")
    await run_cmd([
        "ffmpeg", "-y", "-loop", "1", "-i", image_path, "-i", audio_path,
        "-vf", "scale=1920:1080", "-c:v", "libx264", "-tune", "stillimage",
        "-c:a", "aac", "-b:a", "192k", "-pix_fmt", "yuv420p", "-shortest",
        output_path
    ])
    return output_path

async def merge_video_segments(video_segments: list[str], output_path: str):
    list_file_path = output_path.replace('.mp4', '_list.txt')
    with open(list_file_path, "w") as f:
        for seg in video_segments:
            f.write(f"file '{seg}'\n")

    await run_cmd([
        "ffmpeg", "-y", "-f", "concat", "-safe", "0", "-i", list_file_path,
        "-c", "copy", output_path
    ])
    os.remove(list_file_path)


@app.get("/api/ppt-to-video/download/{task_id}")
async def download_ppt_video(task_id: str):
    if task_id not in ppt_video_progress_store:
        raise HTTPException(status_code=404, detail="Task not found")

    progress = ppt_video_progress_store[task_id]
    if progress["status"] != "completed":
        raise HTTPException(status_code=404, detail="Presentation isn't available")

    video_path = pjoin(progress["task_dir"], "output.mp4")
    if not os.path.exists(video_path):
        raise HTTPException(status_code=404, detail="Presentation not found")

    return FileResponse(
        video_path,
        media_type="video/mp4",
        headers={"Content-Disposition": "attachment; filename=ppt_video.mp4"}
    )


async def ppt_gen(task_id: str, rerun=False):
    if DEBUG:
        importlib.reload(induct)
        importlib.reload(pptgen)
    if rerun:
        task_id = task_id.replace("|", "/")
        active_connections[task_id] = None
        progress_store[task_id] = json.load(open(pjoin(RUNS_DIR, task_id, "task.json")))

    for _ in range(100):
        if task_id in active_connections:
            break
        await asyncio.sleep(0.02)
    else:
        progress_store.pop(task_id)
        return

    task = progress_store.pop(task_id)
    is_topic_mode = bool(task.get("topic"))
    pptx_md5 = task["pptx"]
    pdf_md5 = task.get("pdf")
    generation_config = Config(pjoin(RUNS_DIR, task_id))
    pptx_config = Config(pjoin(RUNS_DIR, "pptx", pptx_md5))
    json.dump(task, open(pjoin(generation_config.RUN_DIR, "task.json"), "w"))
    progress = ProgressManager(
        task_id, TOPIC_STAGES if is_topic_mode else DOCUMENT_STAGES
    )
    if is_topic_mode:
        topic_md5 = hashlib.md5(task["topic"].encode("utf-8")).hexdigest()
        parsedpdf_dir = pjoin(RUNS_DIR, "topic", topic_md5)
    else:
        parsedpdf_dir = pjoin(RUNS_DIR, "pdf", pdf_md5)
    os.makedirs(parsedpdf_dir, exist_ok=True)
    ppt_image_folder = pjoin(pptx_config.RUN_DIR, "slide_images")

    await send_progress(
        active_connections[task_id], "task initialized successfully", 10
    )

    try:
        # ppt parsing
        presentation = Presentation.from_file(
            pjoin(pptx_config.RUN_DIR, "source.pptx"), pptx_config
        )
        if not os.path.exists(ppt_image_folder) or len(
                os.listdir(ppt_image_folder)
        ) != len(presentation):
            await ppt_to_images_async(
                pjoin(pptx_config.RUN_DIR, "source.pptx"), ppt_image_folder
            )
            assert len(os.listdir(ppt_image_folder)) == len(presentation) + len(
                presentation.error_history
            ), "Number of parsed slides and images do not match"

            for err_idx, _ in presentation.error_history:
                os.remove(pjoin(ppt_image_folder, f"slide_{err_idx:04d}.jpg"))
            for i, slide in enumerate(presentation.slides, 1):
                slide.slide_idx = i
                os.rename(
                    pjoin(ppt_image_folder, f"slide_{slide.real_idx:04d}.jpg"),
                    pjoin(ppt_image_folder, f"slide_{slide.slide_idx:04d}.jpg"),
                )

        labler = ImageLabler(presentation, pptx_config)
        if os.path.exists(pjoin(pptx_config.RUN_DIR, "image_stats.json")):
            image_stats = json.load(
                open(pjoin(pptx_config.RUN_DIR, "image_stats.json"), encoding="utf-8")
            )
            labler.apply_stats(image_stats)
        else:
            await labler.caption_images_async(models.vision_model)
            json.dump(
                labler.image_stats,
                open(
                    pjoin(pptx_config.RUN_DIR, "image_stats.json"),
                    "w",
                    encoding="utf-8",
                ),
                ensure_ascii=False,
                indent=4,
            )
        await progress.report_progress()

        if is_topic_mode:
            adapter = DeepResearchAdapter()
            dossier_path = pjoin(parsedpdf_dir, "research_dossier.json")
            if os.path.exists(dossier_path):
                dossier = ResearchDossier.from_dict(
                    json.load(open(dossier_path, encoding="utf-8"))
                )
            else:
                dossier = await run_blocking(
                    adapter.topic_to_dossier,
                    task["topic"],
                    parsedpdf_dir,
                )
                json.dump(
                    dossier.to_dict(),
                    open(dossier_path, "w", encoding="utf-8"),
                    ensure_ascii=False,
                    indent=4,
                )
            await progress.report_progress()

            refined_doc_path = pjoin(parsedpdf_dir, "refined_doc.json")
            if not os.path.exists(refined_doc_path):
                source_doc = await run_blocking(
                    adapter.dossier_to_document,
                    dossier,
                    parsedpdf_dir,
                )
                json.dump(
                    source_doc.to_dict(),
                    open(refined_doc_path, "w", encoding="utf-8"),
                    ensure_ascii=False,
                    indent=4,
                )
            else:
                source_doc = Document.from_dict(
                    json.load(open(refined_doc_path, encoding="utf-8")),
                    parsedpdf_dir,
                    False,
                )
            await progress.report_progress()
        else:
            # pdf parsing
            if not os.path.exists(pjoin(parsedpdf_dir, "source.md")):
                text_content = parse_pdf(
                    pjoin(RUNS_DIR, "pdf", pdf_md5, "source.pdf"),
                    parsedpdf_dir,
                    models.marker_model,
                )
            else:
                text_content = open(
                    pjoin(parsedpdf_dir, "source.md"), encoding="utf-8"
                ).read()
            await progress.report_progress()

            # external media retrieval on source markdown
            enriched_md_path = pjoin(parsedpdf_dir, "source_enriched.md")
            deepresearch_root = os.getenv("PRESENTAGENT_DEEPRESEARCH_ROOT", "").strip()
            if deepresearch_root:
                if not os.path.exists(enriched_md_path):
                    researcher = DocumentMediaResearcher()
                    sync_language_model = (
                        models.language_model.to_sync()
                        if hasattr(models.language_model, "to_sync")
                        else models.language_model
                    )
                    text_content, _ = await run_blocking(
                        researcher.enrich_markdown,
                        markdown_text=text_content,
                        workspace_dir=parsedpdf_dir,
                        language_model=sync_language_model,
                        deepresearch_root=deepresearch_root,
                        conda_env=DOC_MEDIA_DEEPRESEARCH_CONDA_ENV,
                        conda_executable=DOC_MEDIA_DEEPRESEARCH_CONDA_EXE,
                        max_concepts=DOC_MEDIA_MAX_CONCEPTS,
                        max_wait_seconds=DOC_MEDIA_MAX_WAIT_SECONDS,
                        poll_interval_seconds=DOC_MEDIA_POLL_INTERVAL_SECONDS,
                        progress_callback=lambda message: logger.info("doc-media: %s", message),
                    )
                    with open(enriched_md_path, "w", encoding="utf-8") as f:
                        f.write(text_content)
                else:
                    text_content = open(enriched_md_path, encoding="utf-8").read()
            await progress.report_progress()

            # document refine
            if not os.path.exists(pjoin(parsedpdf_dir, "refined_doc.json")):
                source_doc = await Document.from_markdown_async(
                    text_content,
                    models.language_model,
                    models.vision_model,
                    parsedpdf_dir,
                )
                json.dump(
                    source_doc.to_dict(),
                    open(pjoin(parsedpdf_dir, "refined_doc.json"), "w"),
                    ensure_ascii=False,
                    indent=4,
                )
            else:
                source_doc = json.load(open(pjoin(parsedpdf_dir, "refined_doc.json")))
                source_doc = Document.from_dict(source_doc, parsedpdf_dir)
            await progress.report_progress()

        # Slide Induction
        if not os.path.exists(pjoin(pptx_config.RUN_DIR, "slide_induction.json")):
            deepcopy(presentation).save(
                pjoin(pptx_config.RUN_DIR, "template.pptx"), layout_only=True
            )
            await ppt_to_images_async(
                pjoin(pptx_config.RUN_DIR, "template.pptx"),
                pjoin(pptx_config.RUN_DIR, "template_images"),
            )
            slide_inducter = induct.SlideInducterAsync(
                presentation,
                ppt_image_folder,
                pjoin(pptx_config.RUN_DIR, "template_images"),
                pptx_config,
                models.image_model,
                models.language_model,
                models.vision_model,
            )
            layout_induction = await slide_inducter.layout_induct()
            slide_induction = await slide_inducter.content_induct(layout_induction)
            json.dump(
                slide_induction,
                open(
                    pjoin(pptx_config.RUN_DIR, "slide_induction.json"),
                    "w",
                    encoding="utf-8",
                ),
                ensure_ascii=False,
                indent=4,
            )
        else:
            slide_induction = json.load(
                open(
                    pjoin(pptx_config.RUN_DIR, "slide_induction.json"), encoding="utf-8"
                )
            )
        await progress.report_progress()

        # PPT Generation with PPTAgentAsync
        ppt_agent = pptgen.PPTAgentAsync(
            models.text_model,
            models.language_model,
            models.vision_model,
            error_exit=False,
            retry_times=5,
        )
        ppt_agent.set_reference(
            config=generation_config,
            slide_induction=slide_induction,
            presentation=presentation,
        )

        prs, _ = await ppt_agent.generate_pres(
            source_doc=source_doc,
            num_slides=task["numberOfPages"],
        )
        prs.save(pjoin(generation_config.RUN_DIR, "final.pptx"))
        logger.info(f"{task_id}: generation finished")
        await progress.report_progress()
    except Exception as e:
        await progress.fail_stage(str(e))
        traceback.print_exc()


if __name__ == "__main__":
    import uvicorn

    ip = "0.0.0.0"
    uvicorn.run(app, host=ip, port=9297)
