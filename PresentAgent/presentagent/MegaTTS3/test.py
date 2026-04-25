#!/usr/bin/env python
# -*- coding: utf-8 -*-


import os
import tempfile
import argparse
from subprocess import call
import subprocess
from pdf2image import convert_from_path
from pptx import Presentation
# from gtts import gTTS


__author__ = ['chaonan99']


## Sometimes ffmpeg is avconv
FFMPEG_NAME = 'ffmpeg'
# FFMPEG_NAME = 'avconv'

import os
from typing import Optional

from tts.infer_cli import MegaTTS3DiTInfer  # adjust import path as needed
from tts.utils.audio_utils.io import save_wav


def get_tts(
    input_wav_path: str,
    input_text: str,
    output_path: str,
    time_step: int = 32,
    p_w: float = 1.6,
    t_w: float = 2.5,
    device: Optional[str] = None,
) -> str:
    """
    Generate TTS audio from an input WAV file and text prompt.

    Args:
        input_wav_path: Path to the input WAV (prompt) file.
        input_text: Text to synthesize.
        output_path: Path to the output audio file.
        time_step: Diffusion inference steps.
        p_w: Intelligibility weight.
        t_w: Similarity weight.
        device: Device specifier (e.g., 'cuda' or 'cpu'). If None, auto-selected.

    Returns:
        The full path to the generated WAV file.
    """
    # Initialize the inference model
    infer = MegaTTS3DiTInfer(device=device)

    # Read prompt audio
    with open(input_wav_path, 'rb') as f:
        audio_bytes = f.read()

    # Locate corresponding latent file if available
    latent_file = None
    potential_npy = os.path.splitext(input_wav_path)[0] + '.npy'
    if os.path.isfile(potential_npy):
        latent_file = potential_npy

    # Preprocess: extract features and durations
    resource_context = infer.preprocess(audio_bytes, latent_file=latent_file)

    # Synthesize speech
    wav_bytes = infer.forward(
        resource_context,
        input_text,
        time_step=time_step,
        p_w=p_w,
        t_w=t_w
    )

    # Ensure output directory exists and save
    save_wav(wav_bytes, output_path)

    return output_path





def ppt_presenter(pptx_path):
    cmd = ['libreoffice', '--headless', '--convert-to', 'pdf', pptx_path, '--outdir', os.path.dirname(pptx_path)]
    result = subprocess.run(cmd, capture_output=True, text=True)

    pdf_path = os.path.splitext(pptx_path)[0] + '.pdf'
    output_path = os.path.splitext(pptx_path)[0] + '.mp4'
    with tempfile.TemporaryDirectory() as temp_path:
        images_from_path = convert_from_path(pdf_path)
        prs = Presentation(pptx_path)
        assert len(images_from_path) == len(prs.slides)
        for i, (slide, image) in enumerate(zip(prs.slides, images_from_path)):
            if slide.has_notes_slide:
                notes = slide.notes_slide.notes_text_frame.text

                # tts = gTTS(text=notes, lang='en')
                image_path = os.path.join(temp_path, 'frame_{}.jpg'.format(i))
                audio_path = os.path.join(temp_path, 'frame_{}.mp3'.format(i))

                image.save(image_path)
                get_tts("assets/English_prompt.wav", notes, audio_path)
                # tts.save(audio_path)

                ffmpeg_call(image_path, audio_path, temp_path, i)

        video_list = [os.path.join(temp_path, 'frame_{}.ts'.format(i)) \
                      for i in range(len(images_from_path))]
        video_list_str = 'concat:' + '|'.join(video_list)
        ffmpeg_concat(video_list_str, output_path)


def ffmpeg_call(image_path, audio_path, temp_path, i):
    out_path_mp4 = os.path.join(temp_path, 'frame_{}.mp4'.format(i))
    out_path_ts = os.path.join(temp_path, 'frame_{}.ts'.format(i))
    call([FFMPEG_NAME, '-loop', '1', '-y', '-i', image_path, '-i', audio_path,
          '-vf', 'scale=2666:1500', '-c:v', 'libx264', '-tune', 'stillimage', '-c:a', 'aac',
          '-b:a', '192k', '-pix_fmt', 'yuv420p', '-shortest', out_path_mp4])

    call([FFMPEG_NAME, '-y', '-i', out_path_mp4, '-c', 'copy',
          '-bsf:v', 'h264_mp4toannexb', '-f', 'mpegts', out_path_ts])


def ffmpeg_concat(video_list_str, out_path):
    call([FFMPEG_NAME, '-y', '-f', 'mpegts', '-i', '{}'.format(video_list_str),
          '-c', 'copy', '-bsf:a', 'aac_adtstoasc', out_path])


def main():
    parser = argparse.ArgumentParser(description='PPT Presenter help.')
    parser.add_argument('--pptx', default='../../ppagent_2025-06-29_152592d9-df14-48d0-b6de-99fa7fe4fdac.pptx', help='input pptx path')
    args = parser.parse_args()
    ppt_presenter(args.pptx)


if __name__ == '__main__':
    main()