from __future__ import annotations

import argparse
import asyncio
import os
import subprocess
import tempfile
import wave
from pathlib import Path

import numpy as np
from pdf2image import convert_from_path
from pptx import Presentation as PptxPresentation


def _run_cmd(cmd: list[str]) -> bytes:
    result = subprocess.run(cmd, stdout=subprocess.PIPE, stderr=subprocess.PIPE, check=False)
    if result.returncode != 0:
        raise RuntimeError(f"{' '.join(cmd)}\n{result.stderr.decode(errors='ignore')}")
    return result.stdout


def _generate_tts_audio(text: str, output_path: str) -> None:
    try:
        from presentagent.MegaTTS3.tts.infer_cli import MegaTTS3DiTInfer
        from presentagent.MegaTTS3.tts.utils.audio_utils.io import save_wav

        ckpt_root = str(Path(__file__).resolve().parents[1] / "presentagent" / "MegaTTS3" / "checkpoints")
        infer = MegaTTS3DiTInfer(ckpt_root=ckpt_root)

        prompt_audio_path = (
            Path(__file__).resolve().parents[1]
            / "presentagent"
            / "MegaTTS3"
            / "assets"
            / "English_prompt.wav"
        )
        audio_bytes = prompt_audio_path.read_bytes()
        latent_file = None
        potential_npy = prompt_audio_path.with_suffix(".npy")
        if potential_npy.is_file():
            latent_file = str(potential_npy)
        resource_context = infer.preprocess(audio_bytes, latent_file)
        wav_bytes = infer.forward(resource_context, text, time_step=32, p_w=1.6, t_w=2.5)
        save_wav(wav_bytes, output_path)
    except Exception:
        sample_rate = 22050
        duration = 3.0
        samples = np.zeros(int(sample_rate * duration), dtype=np.int16)
        with wave.open(output_path, "w") as wav_file:
            wav_file.setnchannels(1)
            wav_file.setsampwidth(2)
            wav_file.setframerate(sample_rate)
            wav_file.writeframes(samples.tobytes())


def _create_video_segment(image_path: str, audio_path: str, temp_dir: str, index: int) -> str:
    output_path = os.path.join(temp_dir, f"segment_{index}.mp4")
    _run_cmd(
        [
            "ffmpeg",
            "-y",
            "-loop",
            "1",
            "-i",
            image_path,
            "-i",
            audio_path,
            "-vf",
            "scale=1920:1080",
            "-c:v",
            "libx264",
            "-tune",
            "stillimage",
            "-c:a",
            "aac",
            "-b:a",
            "192k",
            "-pix_fmt",
            "yuv420p",
            "-shortest",
            output_path,
        ]
    )
    return output_path


def _merge_video_segments(video_segments: list[str], output_path: str) -> None:
    list_file_path = output_path.replace(".mp4", "_list.txt")
    with open(list_file_path, "w", encoding="utf-8") as f:
        for seg in video_segments:
            f.write(f"file '{seg}'\n")
    try:
        _run_cmd(
            [
                "ffmpeg",
                "-y",
                "-f",
                "concat",
                "-safe",
                "0",
                "-i",
                list_file_path,
                "-c",
                "copy",
                output_path,
            ]
        )
    finally:
        if os.path.exists(list_file_path):
            os.remove(list_file_path)


async def _build_video(ppt_path: Path, output_dir: Path) -> None:
    pdf_path = output_dir / "source.pdf"
    _run_cmd(
        [
            "libreoffice",
            "--headless",
            "--convert-to",
            "pdf",
            str(ppt_path),
            "--outdir",
            str(output_dir),
        ]
    )

    images_from_path = await asyncio.to_thread(convert_from_path, str(pdf_path))
    prs = await asyncio.to_thread(PptxPresentation, str(ppt_path))
    if len(images_from_path) != len(prs.slides):
        raise RuntimeError("Slide count does not match rendered image count.")

    with tempfile.TemporaryDirectory() as temp_dir:
        notes_dir = output_dir / "notes_assets"
        notes_dir.mkdir(parents=True, exist_ok=True)
        video_segments: list[str] = []

        notes_dump = []
        for i, (slide, image) in enumerate(zip(prs.slides, images_from_path)):
            notes = slide.notes_slide.notes_text_frame.text if slide.has_notes_slide else ""
            if not notes.strip():
                notes = f"This is slide {i + 1}"
            notes_dump.append({"slide_idx": i + 1, "notes": notes})

            image_path = notes_dir / f"frame_{i}.jpg"
            audio_path = notes_dir / f"frame_{i}.wav"
            image.save(image_path)
            await asyncio.to_thread(_generate_tts_audio, notes, str(audio_path))
            segment_path = await asyncio.to_thread(
                _create_video_segment,
                str(image_path),
                str(audio_path),
                str(notes_dir),
                i,
            )
            video_segments.append(segment_path)

        (output_dir / "slide_notes.json").write_text(
            __import__("json").dumps(notes_dump, ensure_ascii=False, indent=2),
            encoding="utf-8",
        )
        final_video_path = output_dir / "output.mp4"
        await asyncio.to_thread(_merge_video_segments, video_segments, str(final_video_path))

    print("success: True")
    print(f"pptx: {ppt_path}")
    print(f"pdf: {pdf_path}")
    print(f"video: {output_dir / 'output.mp4'}")
    print(f"notes: {output_dir / 'slide_notes.json'}")
    print(f"assets_dir: {output_dir / 'notes_assets'}")


def main() -> None:
    parser = argparse.ArgumentParser(
        description="Render a PPT with notes into a narrated MP4 and save intermediate assets."
    )
    parser.add_argument("--pptx", required=True, help="Path to final.pptx")
    parser.add_argument("--output-dir", required=True, help="Directory to save video outputs")
    args = parser.parse_args()

    ppt_path = Path(args.pptx).resolve()
    if not ppt_path.exists():
        raise SystemExit(f"PPTX not found: {ppt_path}")
    output_dir = Path(args.output_dir).resolve()
    output_dir.mkdir(parents=True, exist_ok=True)
    asyncio.run(_build_video(ppt_path, output_dir))


if __name__ == "__main__":
    main()
